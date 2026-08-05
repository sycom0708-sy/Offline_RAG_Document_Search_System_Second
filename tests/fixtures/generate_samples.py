"""테스트용 샘플 문서 생성 (T1.13).

바이너리 샘플을 저장소에 커밋하는 대신 코드로 생성한다.
모든 샘플은 한국어 본문 + 표 1개 + 이미지 1개를 포함해 표/이미지 분리 추출을 검증할 수 있게 한다.
"""

from __future__ import annotations

import io
import zipfile
from pathlib import Path

BODY_TEXT = "오프라인 문서 검색 시스템 개요"
BODY_PARAGRAPH = "사내 문서를 완전 오프라인 환경에서 검색하기 위한 시스템이다."
TABLE_HEADER = ["구분", "최소 사양", "권장 사양"]
TABLE_ROWS = [
    ["RAM", "8GB", "16GB"],
    ["GPU", "없음", "없음"],
]


def make_png_bytes(width: int = 40, height: int = 30, color: tuple[int, int, int] = (200, 60, 60)) -> bytes:
    from PIL import Image

    buffer = io.BytesIO()
    Image.new("RGB", (width, height), color).save(buffer, format="PNG")
    return buffer.getvalue()


def make_txt(path: Path, encoding: str = "utf-8") -> Path:
    content = f"{BODY_TEXT}\n{BODY_PARAGRAPH}\n한글 인코딩 테스트: 가나다라마바사"
    path.write_bytes(content.encode(encoding))
    return path


def make_docx(path: Path) -> Path:
    import docx

    document = docx.Document()
    document.add_heading(BODY_TEXT, level=1)
    document.add_paragraph(BODY_PARAGRAPH)

    table = document.add_table(rows=len(TABLE_ROWS) + 1, cols=len(TABLE_HEADER))
    for col, name in enumerate(TABLE_HEADER):
        table.cell(0, col).text = name
    for row_index, row in enumerate(TABLE_ROWS, start=1):
        for col, value in enumerate(row):
            table.cell(row_index, col).text = value

    document.add_paragraph("표 아래 문단입니다.")
    document.add_picture(io.BytesIO(make_png_bytes()))
    document.save(str(path))
    return path


def make_xlsx(path: Path) -> Path:
    import openpyxl
    from openpyxl.drawing.image import Image as XlsxImage

    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "사양표"
    sheet.append(TABLE_HEADER)
    for row in TABLE_ROWS:
        sheet.append(row)

    image_path = path.parent / "_xlsx_embedded.png"
    image_path.write_bytes(make_png_bytes())
    sheet.add_image(XlsxImage(str(image_path)), "E2")

    second = workbook.create_sheet("메모")
    second.append([BODY_TEXT])
    second.append([BODY_PARAGRAPH])

    workbook.save(str(path))
    return path


def make_pptx(path: Path) -> Path:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = BODY_TEXT
    slide.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(1)).text_frame.text = BODY_PARAGRAPH

    table_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = table_slide.shapes.add_table(
        len(TABLE_ROWS) + 1, len(TABLE_HEADER), Inches(1), Inches(1), Inches(6), Inches(2)
    )
    for col, name in enumerate(TABLE_HEADER):
        shape.table.cell(0, col).text = name
    for row_index, row in enumerate(TABLE_ROWS, start=1):
        for col, value in enumerate(row):
            shape.table.cell(row_index, col).text = value

    table_slide.shapes.add_picture(io.BytesIO(make_png_bytes()), Inches(7), Inches(1))
    presentation.save(str(path))
    return path


def make_pdf(path: Path) -> Path:
    import pymupdf

    document = pymupdf.open()
    page = document.new_page()
    # 내장 한국어 폰트("korea")를 써야 한글 글리프가 매핑돼 텍스트 추출이 가능하다.
    page.insert_text((60, 80), BODY_TEXT, fontname="korea", fontsize=16)
    page.insert_text((60, 110), BODY_PARAGRAPH, fontname="korea", fontsize=11)

    # 괘선을 직접 그려 find_tables()가 인식할 수 있는 표를 만든다.
    left, top, col_w, row_h = 60, 150, 140, 24
    grid = [TABLE_HEADER] + TABLE_ROWS
    for row_index, row in enumerate(grid):
        for col_index, value in enumerate(row):
            rect = pymupdf.Rect(
                left + col_index * col_w,
                top + row_index * row_h,
                left + (col_index + 1) * col_w,
                top + (row_index + 1) * row_h,
            )
            page.draw_rect(rect, color=(0, 0, 0), width=0.8)
            page.insert_text((rect.x0 + 5, rect.y0 + 16), value, fontname="korea", fontsize=10)

    page.insert_image(pymupdf.Rect(60, 260, 160, 335), stream=make_png_bytes())

    diagram_page = document.new_page()
    diagram_page.insert_text((60, 60), "도면 페이지", fontname="korea", fontsize=12)
    for index in range(6):
        diagram_page.draw_circle(
            pymupdf.Point(120 + index * 60, 300), 40, color=(0, 0, 1), width=2
        )

    document.save(str(path))
    document.close()
    return path


def make_hwpx(path: Path) -> Path:
    """OWPML 구조의 최소 HWPX를 직접 조립한다 (한/글 없이 생성하기 위함)."""
    ns = "http://www.hancom.co.kr/hwpml/2011/paragraph"
    rows_xml = "".join(
        "<hp:tr>"
        + "".join(
            f'<hp:tc><hp:subList><hp:p><hp:run><hp:t>{cell}</hp:t></hp:run></hp:p></hp:subList></hp:tc>'
            for cell in row
        )
        + "</hp:tr>"
        for row in [TABLE_HEADER] + TABLE_ROWS
    )
    section = (
        f'<?xml version="1.0" encoding="UTF-8"?>'
        f'<hs:sec xmlns:hp="{ns}" xmlns:hs="http://www.hancom.co.kr/hwpml/2011/section">'
        f"<hp:p><hp:run><hp:t>{BODY_TEXT}</hp:t></hp:run></hp:p>"
        f"<hp:p><hp:run><hp:t>{BODY_PARAGRAPH}</hp:t></hp:run></hp:p>"
        f"<hp:tbl>{rows_xml}</hp:tbl>"
        f"<hp:p><hp:run><hp:t>표 아래 문단입니다.</hp:t></hp:run></hp:p>"
        f"</hs:sec>"
    )

    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("mimetype", "application/hwp+zip")
        archive.writestr("Contents/section0.xml", section)
        archive.writestr("BinData/image1.png", make_png_bytes())
    return path


GENERATORS = {
    "sample.txt": make_txt,
    "sample.docx": make_docx,
    "sample.xlsx": make_xlsx,
    "sample.pptx": make_pptx,
    "sample.pdf": make_pdf,
    "sample.hwpx": make_hwpx,
}


def generate_all(output_dir: Path) -> dict[str, Path]:
    output_dir.mkdir(parents=True, exist_ok=True)
    return {name: generator(output_dir / name) for name, generator in GENERATORS.items()}


if __name__ == "__main__":
    target = Path(__file__).resolve().parents[1] / "samples" / "generated"
    for name, created in generate_all(target).items():
        print(f"{name}: {created}")
