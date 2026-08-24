"""PPTX 파서 (T1.6) — python-pptx 기반 슬라이드별 텍스트/표/이미지 추출."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from parser.base import BaseParser, DocumentReadError
from parser.schema import ImageData, ParsedDocument, TableData
from parser.utils.headings import MAX_HEADING_CHARS, pick_largest_top_line
from parser.utils.libreoffice import LibreOfficeError
from parser.utils.render import render_pages

# 제목 길이 상한 (T10.31) — PDF·docx 쪽과 같은 기준.


class PptxParser(BaseParser):
    extensions = (".pptx",)

    def __init__(self, asset_dir: str | Path | None = None, capture_vector_shapes: bool = False) -> None:
        super().__init__(asset_dir)
        self._capture_vector_shapes = capture_vector_shapes

    def _parse(self, path: Path, document: ParsedDocument) -> None:
        try:
            presentation = Presentation(str(path))
        except Exception as exc:
            raise DocumentReadError(f"PPTX를 열 수 없습니다: {path.name} ({exc})") from exc

        asset_dir = self.asset_dir_for(path)
        first_title = ""

        for slide_index, slide in enumerate(presentation.slides, start=1):
            texts: list[str] = []
            # 슬라이드 제목은 표·이미지보다 뒤에 나올 수도 있어 먼저 읽어 둔다
            # (T10.31) — 도형 순회 중에 채우면 앞쪽 청크가 제목을 못 받는다.
            slide_heading = self._slide_title(slide)
            for shape in self._iter_shapes(slide.shapes):
                if shape.has_table:
                    table_data = self._read_table(shape.table)
                    if table_data is not None:
                        document.chunks.append(
                            self.make_table_chunk(
                                document, table_data, page_or_slide=slide_index,
                                heading=slide_heading,
                            )
                        )
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    self._extract_image(document, shape, slide_index, asset_dir, slide_heading)
                elif shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        texts.append(text)
                        if not first_title:
                            first_title = text.splitlines()[0][:200]

            body = "\n".join(texts).strip()
            if body:
                document.chunks.append(
                    self.make_text_chunk(
                        document, body, page_or_slide=slide_index, heading=slide_heading
                    )
                )

        document.title = (presentation.core_properties.title or "").strip() or first_title or path.stem
        self._capture_drawings(document, path, asset_dir)

    @staticmethod
    def _slide_title(slide) -> str:
        """슬라이드의 제목 — 플레이스홀더 우선(T10.31), 없으면 글꼴+위치 폴백(T10.33).

        pptx는 제목을 별도 플레이스홀더로 들고 있어 있으면 추측이 필요 없다 —
        PDF가 글꼴 크기로, docx가 스타일 이름으로 푸는 것을 여기서는 정확히
        알 수 있다. 다만 이 코퍼스는 슬라이드의 상당수가 제목을 플레이스홀더가
        아니라 일반 텍스트 상자로 쓴다(`_font_size_heading` 참고).
        """
        try:
            title_shape = slide.shapes.title
        except Exception:
            title_shape = None
        if title_shape is not None and title_shape.has_text_frame:
            title = title_shape.text_frame.text.strip()
            if title:
                return title.splitlines()[0][:MAX_HEADING_CHARS]
        return PptxParser._font_size_heading(slide)

    @classmethod
    def _font_size_heading(cls, slide) -> str:
        """제목 플레이스홀더가 없는 슬라이드를 위한 글꼴 크기+위치 폴백 (T10.33).

        `pick_largest_top_line()`이 "가장 크면서 가장 위"인 도형만 제목으로
        인정한다 — 판단 근거는 그 함수 docstring 참고.
        """
        sized: list[tuple[float, float, str]] = []
        for shape in cls._iter_shapes(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            sizes = [
                run.font.size.pt
                for paragraph in shape.text_frame.paragraphs
                for run in paragraph.runs
                if run.font.size is not None
            ]
            if not sizes:
                continue
            top = shape.top if shape.top is not None else float("inf")
            sized.append((max(sizes), top, text.splitlines()[0]))
        return pick_largest_top_line(sized)

    @staticmethod
    def _iter_shapes(shapes):
        """그룹 도형 안쪽까지 펼쳐서 순회한다."""
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from PptxParser._iter_shapes(shape.shapes)
            else:
                yield shape

    @staticmethod
    def _read_table(table) -> TableData | None:
        rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        return TableData.from_rows(rows)

    def _extract_image(
        self, document: ParsedDocument, shape, slide_index: int, asset_dir: Path,
        heading: str = "",
    ) -> None:
        try:
            image = shape.image
            blob = image.blob
            ext = image.ext
        except Exception as exc:
            self.record_error(document, f"{slide_index}번 슬라이드 이미지 추출 실패: {exc}")
            return

        out_path = asset_dir / f"slide{slide_index:04d}_{shape.shape_id}.{ext}"
        out_path.write_bytes(blob)
        document.chunks.append(
            self.make_image_chunk(
                document,
                ImageData(
                    image_path=str(out_path),
                    caption=(shape.name or "").strip(),
                    origin="extracted",
                ),
                page_or_slide=slide_index,
                heading=heading,
            )
        )

    def _capture_drawings(self, document: ParsedDocument, path: Path, asset_dir: Path) -> None:
        if not self._capture_vector_shapes:
            return
        try:
            captured = render_pages(path, asset_dir)
        except LibreOfficeError as exc:
            self.record_error(document, f"벡터 도형 캡처 실패: {exc}")
            return

        for slide_no, image_path in sorted(captured.items()):
            document.chunks.append(
                self.make_image_chunk(
                    document,
                    ImageData(
                        image_path=str(image_path),
                        caption=f"{slide_no}번 슬라이드 캡처",
                        origin="rendered",
                    ),
                    page_or_slide=slide_no,
                )
            )
